from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class PPTManager:
    """Toolkit for generating standardized and dynamic PowerPoint presentations."""

    # Default Theme Colors
    THEMES = {
        "modern": {
            "bg": RGBColor(248, 249, 250),
            "text_main": RGBColor(17, 24, 39),
            "text_body": RGBColor(75, 85, 99),
            "accent": RGBColor(237, 234, 229),
            "white": RGBColor(255, 255, 255)
        },
        "classic_blue": {
            "bg": RGBColor(255, 255, 255),
            "text_main": RGBColor(12, 77, 162),
            "text_body": RGBColor(80, 80, 80),
            "accent": RGBColor(12, 77, 162),
            "white": RGBColor(255, 255, 255)
        }
    }

    def __init__(self, theme_name="modern", font_name="맑은 고딕", ratio="16:9"):
        self.prs = Presentation()
        if ratio == "16:9":
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(5.625)
        
        self.theme = self.THEMES.get(theme_name, self.THEMES["modern"])
        self.font_name = font_name

    def _apply_background(self, slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(self, title_text, subtitle_text):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        self._apply_background(slide, self.theme["bg"])
        
        title = slide.shapes.title
        title.text = title_text
        for p in title.text_frame.paragraphs:
            p.font.color.rgb = self.theme["text_main"]
            p.font.name = self.font_name
            p.font.bold = True
            p.font.size = Pt(40)
            p.alignment = PP_ALIGN.LEFT
            
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        for p in subtitle.text_frame.paragraphs:
            p.font.color.rgb = self.theme["text_body"]
            p.font.name = self.font_name
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.LEFT

    def add_content_slide(self, title_text, items):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5]) # blank with title
        self._apply_background(slide, self.theme["white"])
        
        title = slide.shapes.title
        title.text = title_text
        for p in title.text_frame.paragraphs:
            p.font.color.rgb = self.theme["text_main"]
            p.font.name = self.font_name
            p.font.bold = True
            p.font.size = Pt(28)
            p.alignment = PP_ALIGN.LEFT

        left = Inches(0.8); top = Inches(1.6); width = Inches(8.4); height = Inches(3.4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(16)
            p.font.name = self.font_name
            p.font.color.rgb = self.theme["text_body"]
            p.space_after = Pt(14)

    def add_table_slide(self, title_text, headers, rows):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._apply_background(slide, self.theme["white"])
        
        title = slide.shapes.title
        title.text = title_text
        for p in title.text_frame.paragraphs:
            p.font.color.rgb = self.theme["text_main"]
            p.font.name = self.font_name
            p.font.bold = True
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.LEFT
            
        rows_count = len(rows) + 1
        cols_count = len(headers)
        
        left = Inches(0.5) if cols_count < 4 else Inches(0.2)
        width = Inches(9.0) if cols_count < 4 else Inches(9.6)
        top = Inches(1.4)
        height = Inches(1.0)
        
        table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
        table = table_shape.table
        
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme["bg"]
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = self.theme["text_main"]
                p.font.name = self.font_name
                p.font.bold = True
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                
        for r_idx, row in enumerate(rows):
            for c_idx, cell_data in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(cell_data)
                for p in cell.text_frame.paragraphs:
                    p.font.name = self.font_name
                    p.font.color.rgb = self.theme["text_body"]
                    p.font.size = Pt(11)
                    p.alignment = PP_ALIGN.CENTER
                    if c_idx == 0:
                        p.font.color.rgb = self.theme["text_main"]
                        p.font.bold = True

    def save(self, filepath):
        self.prs.save(filepath)
        print(f"Presentation saved to {filepath}")

if __name__ == '__main__':
    print("PPTManager loaded. Instantiate to create standard, compliant presentations.")
